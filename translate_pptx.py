from pptx import Presentation
from pptx.util import Pt
from googletrans import Translator
import os
import time
import io
from PIL import Image

def translate_pptx(input_path, output_path):
    # Load the presentation
    presentation = Presentation(input_path)

    # Initialize translator
    translator = Translator()

    # Create output directory for images
    image_output_dir = os.path.join(os.path.dirname(output_path), 'pptx_images')
    os.makedirs(image_output_dir, exist_ok=True)

    # Function to copy font properties
    def copy_font_properties(original_run, new_run):
        new_run.font.name = original_run.font.name
        new_run.font.size = original_run.font.size
        new_run.font.bold = original_run.font.bold
        new_run.font.italic = original_run.font.italic
        if original_run.font.color:
            try:
                if original_run.font.color.rgb:
                    new_run.font.color.rgb = original_run.font.color.rgb
            except AttributeError:
                pass
            try:
                if original_run.font.color.theme_color:
                    new_run.font.color.theme_color = original_run.font.color.theme_color
            except AttributeError:
                pass

    # Function to translate text with retry mechanism
    def translate_text(text, retries=3):
        for attempt in range(retries):
            try:
                if text.strip():  # Only attempt to translate non-empty text
                    translated = translator.translate(text, dest='en').text
                    return translated
                else:
                    return text  # Return original text if empty
            except Exception as e:
                if attempt < retries - 1:
                    time.sleep(2)  # Wait for 2 seconds before retrying
                else:
                    print(f"Translation error: {e}")
                    return text  # Return original text in case of repeated failure

    # Function to handle translation and text placement
    def handle_translation(paragraph):
        original_runs = paragraph.runs
        paragraph.text = ""

        # Add original text runs back
        for run in original_runs:
            new_run = paragraph.add_run()
            new_run.text = run.text
            copy_font_properties(run, new_run)

        # Collect translated texts
        translated_texts = []
        for run in original_runs:
            translated_text = translate_text(run.text)
            if translated_text.strip():
                translated_texts.append(translated_text)

        # Add translated text as a single run
        if translated_texts:
            translated_run = paragraph.add_run()
            translated_run.text = '\n' + " ".join(translated_texts)
            copy_font_properties(original_runs[0], translated_run)
            translated_run.font.size = Pt(10)  # Adjust font size for translated text
            translated_run.font.italic = True  # Example style change: make italic

    # Iterate through slides
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                paragraphs = list(text_frame.paragraphs)  # Create a list to avoid modifying the original iterable
                for paragraph in paragraphs:
                    handle_translation(paragraph)

            # Check for images
            if shape.shape_type == 13:  # Picture shape
                try:
                    image = shape.image
                    image_bytes = io.BytesIO(image.blob)
                    img = Image.open(image_bytes)
                    # Save image with unique name
                    img_filename = f"image_slide_{slide.slide_id}_shape_{shape.shape_id}.png"
                    img.save(os.path.join(image_output_dir, img_filename))
                except Exception as e:
                    print(f"Image processing error: {e}")

            # Check for other shapes like arrows
            if shape.shape_type in [6, 7]:  # Arrow shapes
                try:
                    left = shape.left
                    top = shape.top
                    shape.left = left
                    shape.top = top
                except Exception as e:
                    print(f"Arrow positioning error: {e}")

    # Save the modified presentation
    try:
        presentation.save(output_path)
    except Exception as e:
        print(f"Error saving presentation: {e}")

if __name__ == "__main__":
    input_path = '/mnt/data/Translated_Networking.pptx'
    output_path = '/mnt/data/Translated_Networking_Processed.pptx'
    translate_pptx(input_path, output_path)
