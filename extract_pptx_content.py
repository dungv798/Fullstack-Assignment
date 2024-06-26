import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from googletrans import Translator

def extract_translate_append_pptx(input_pptx_path, output_pptx_path):
    translator = Translator()
    presentation = Presentation(input_pptx_path)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                original_text = paragraph.text
                if original_text.strip():  # Only translate non-empty text
                    translated_text = translator.translate(original_text, dest='en').text

                    # Add translated text
                    new_paragraph = shape.text_frame.add_paragraph()
                    new_paragraph.text = translated_text
                    new_paragraph.font.size = paragraph.font.size
                    new_paragraph.font.name = paragraph.font.name
                    new_paragraph.font.bold = paragraph.font.bold
                    new_paragraph.font.italic = paragraph.font.italic

                    # Set font color only if it is available
                    if isinstance(paragraph.font.color._color, RGBColor):
                        new_paragraph.font.color.rgb = paragraph.font.color.rgb
                        print(f"Original color: {paragraph.font.color.rgb}")
                        print(f"New color: {new_paragraph.font.color.rgb}")
                    else:
                        print("Original color: None, defaulting to black")
                        new_paragraph.font.color.rgb = RGBColor(0, 0, 0)
                        print(f"New color: {new_paragraph.font.color.rgb}")

    presentation.save(output_pptx_path)

if __name__ == "__main__":
    input_pptx_path = "path/to/input.pptx"  # Replace with your input file path
    output_dir = "path/to/output/dir"  # Replace with your output directory
    extract_translate_append_pptx(input_pptx_path, os.path.join(output_dir, "translated_pptx.pptx"))
